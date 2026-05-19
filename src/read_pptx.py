from pptx import Presentation
import os

def main():
    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "../data/AI_wellbore_geology_prediction_task_en.pptx"))
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    print(f"Loaded presentation: {pptx_path}")
    print(f"Number of slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text.strip())

if __name__ == "__main__":
    main()
